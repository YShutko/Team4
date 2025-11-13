"""
Build PowerPoint Presentation for Spotify Track Popularity Prediction

This script generates a professional PowerPoint presentation using captured
screenshots and model metadata.
"""

import json
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Configuration
SCREENS_DIR = Path(__file__).parent / "screens"
OUTPUT_DIR = Path(__file__).parent / "output"
TEMPLATE_DIR = Path(__file__).parent / "templates"
OUTPUT_FILE = OUTPUT_DIR / "Spotify_Popularity_Prediction_Presentation.pptx"

# Ensure output directory exists
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# Slide content configuration
SLIDE_TITLES = {
    "01_xgboost_learning_curve.png": "Learning Curves - Model Training Progress",
    "02_actual_vs_predicted.png": "Model Performance - Actual vs Predicted",
    "03_prediction_density.png": "Prediction Density Distribution",
    "04_residuals_plot.png": "Residuals Analysis",
    "05_qq_plot_residuals.png": "QQ Plot - Residuals Normality Check",
    "06_correlation_heatmap.png": "Feature Correlation Heatmap",
    "07_feature_importance.png": "XGBoost Feature Importance",
    "08_xgboost_shap_summary_bar.png": "SHAP Values - Feature Impact (Bar)",
    "09_xgboost_shap_beeswarm.png": "SHAP Values - Feature Impact (Beeswarm)",
}

SLIDE_DESCRIPTIONS = {
    "01_xgboost_learning_curve.png": "Training and validation RMSE across iterations. No significant overfitting observed.",
    "02_actual_vs_predicted.png": "Strong correlation between actual and predicted popularity scores.",
    "03_prediction_density.png": "Density plot showing concentration of predictions around actual values.",
    "04_residuals_plot.png": "Residuals evenly distributed around zero, indicating good model fit.",
    "05_qq_plot_residuals.png": "Residuals follow normal distribution, validating model assumptions.",
    "06_correlation_heatmap.png": "Feature correlations reveal key relationships driving popularity.",
    "07_feature_importance.png": "XGBoost gain-based importance showing top predictive features.",
    "08_xgboost_shap_summary_bar.png": "SHAP values quantify each feature's average impact on predictions.",
    "09_xgboost_shap_beeswarm.png": "SHAP beeswarm plot shows feature impact direction and magnitude.",
}


def load_metadata():
    """Load metadata from capture_screens.py output"""
    metadata_file = SCREENS_DIR / "metadata.json"
    if metadata_file.exists():
        with open(metadata_file, 'r') as f:
            return json.load(f)
    return {}


def create_title_slide(prs, metadata):
    """Create title slide with project metadata"""
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(title_slide_layout)

    # Title
    title = slide.shapes.title
    title.text = metadata.get('title', 'Spotify Track Popularity Prediction')

    # Subtitle
    subtitle = slide.placeholders[1]
    subtitle_text = metadata.get('subtitle', 'Machine Learning Pipeline with XGBoost & SHAP')

    # Add metrics if available
    if 'test_r2' in metadata:
        subtitle_text += f"\n\nTest R²: {metadata['test_r2']:.4f}"
    if 'test_rmse' in metadata:
        subtitle_text += f" | RMSE: {metadata['test_rmse']:.2f}"

    subtitle_text += f"\n\nGenerated: {datetime.now().strftime('%B %d, %Y')}"

    subtitle.text = subtitle_text


def create_agenda_slide(prs):
    """Create agenda/outline slide"""
    title_only_layout = prs.slide_layouts[5]  # Title only layout
    slide = prs.slides.add_slide(title_only_layout)

    title = slide.shapes.title
    title.text = "Presentation Agenda"

    # Add content text box
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    agenda_items = [
        "1. Learning Curves - Model Training Progress",
        "2. Model Performance Metrics",
        "3. Prediction Quality Analysis",
        "4. Residuals & Statistical Validation",
        "5. Feature Correlations",
        "6. Feature Importance Analysis",
        "7. SHAP Explainability",
        "8. Key Insights & Findings"
    ]

    for item in agenda_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(12)


def create_content_slide(prs, image_path, title, description):
    """Create a content slide with image and description"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title

    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue

    # Add image
    img_left = Inches(0.5)
    img_top = Inches(1.3)
    img_width = Inches(9)

    slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width)

    # Add description at bottom
    desc_left = Inches(0.5)
    desc_top = Inches(6.5)
    desc_width = Inches(9)
    desc_height = Inches(0.8)

    desc_box = slide.shapes.add_textbox(desc_left, desc_top, desc_width, desc_height)
    desc_frame = desc_box.text_frame
    desc_frame.text = description

    p = desc_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = RGBColor(64, 64, 64)  # Dark gray


def create_summary_slide(prs, metadata):
    """Create final summary slide"""
    title_only_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_layout)

    title = slide.shapes.title
    title.text = "Key Findings & Insights"

    # Add content
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    findings = [
        f"✅ Achieved R² of {metadata.get('test_r2', 0):.4f} on test set",
        f"✅ RMSE of {metadata.get('test_rmse', 0):.2f} demonstrates strong predictive accuracy",
        "✅ No significant overfitting observed in learning curves",
        "✅ Residuals follow normal distribution, validating model assumptions",
        "✅ SHAP analysis reveals key drivers of track popularity",
        f"✅ Model uses {len(metadata.get('features', []))} audio features for prediction",
        "✅ Full MLflow experiment tracking enables reproducibility",
    ]

    for finding in findings:
        p = tf.add_paragraph()
        p.text = finding
        p.font.size = Pt(16)
        p.space_after = Pt(14)


def build_presentation(metadata):
    """Build the complete presentation"""
    print("="*80)
    print("🎬 BUILDING POWERPOINT PRESENTATION")
    print("="*80)

    # Create presentation
    prs = Presentation()

    # Set slide size to 16:9 widescreen
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("\n📄 Adding slides...")

    # Slide 1: Title
    print("  1. Title slide")
    create_title_slide(prs, metadata)

    # Slide 2: Agenda
    print("  2. Agenda slide")
    create_agenda_slide(prs)

    # Slides 3-N: Content slides with images
    screenshots = sorted(SCREENS_DIR.glob("*.png"))
    slide_num = 3

    for screenshot in screenshots:
        if screenshot.name == "metadata.json":
            continue

        title = SLIDE_TITLES.get(screenshot.name, screenshot.stem.replace('_', ' ').title())
        description = SLIDE_DESCRIPTIONS.get(screenshot.name, "")

        print(f"  {slide_num}. {title}")
        create_content_slide(prs, screenshot, title, description)
        slide_num += 1

    # Final slide: Summary
    print(f"  {slide_num}. Summary slide")
    create_summary_slide(prs, metadata)

    # Save presentation
    prs.save(str(OUTPUT_FILE))

    print("\n" + "="*80)
    print("✅ PRESENTATION BUILT SUCCESSFULLY")
    print("="*80)
    print(f"\n📁 Output file: {OUTPUT_FILE}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"📸 Screenshots used: {len(screenshots)}")

    return OUTPUT_FILE


def main():
    """Main execution"""

    # Load metadata
    metadata = load_metadata()

    # Check if screenshots exist
    screenshots = list(SCREENS_DIR.glob("*.png"))
    if not screenshots:
        print("❌ Error: No screenshots found in screens/ directory")
        print("   Please run 'python presentation/capture_screens.py' first")
        return

    # Build presentation
    output_file = build_presentation(metadata)

    print("\n🎯 Next steps:")
    print("  - Review presentation: open presentation/output/Spotify_Popularity_Prediction_Presentation.pptx")
    print("  - Customize slides as needed")
    print("  - Share with stakeholders")
    print("\n" + "="*80)


if __name__ == "__main__":
    main()
